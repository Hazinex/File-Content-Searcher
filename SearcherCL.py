# Necessary imports for non-native file handling, directory traversal and progress bars
import os
import docx
import openpyxl
import PyPDF2
import pptx
from tqdm import tqdm
import tkinter as tk

# This function makes a tuple into all lowercase
def lowercaseTuple(tuple_val):
    if tuple_val is None:
        return None
    lowercase_vals = []
    for val in tuple_val:
        if val is not None:
            lowercase_vals.append(str(val).lower())
        else:
            lowercase_vals.append(None)
    return tuple(lowercase_vals)

# This function returns the path if it contains the desired keyword 
def readPDF(path, query):
    with open(path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
        if query.lower() in text.casefold() :
            return path

# This function returns the path if it contains the desired keyword 
def readDOCX(path, query):
    document = docx.Document(path)
    for paragraph in document.paragraphs:
        if query.lower() in paragraph.text.casefold():
            return path

# This function returns the path if it contains the desired keyword 
def readPPTX(path, query):
    prs = pptx.Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and query.lower() in shape.text.casefold():
                return path

# This function returns the path if it contains the desired keyword 
def readXLSX(path, query):
    workbook = openpyxl.load_workbook(path)
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            row = lowercaseTuple(row)
            if query.lower() in row:
                return path

print("\n-----------------------------------------------------------------------------------\n")
print("Welcome to the Knowledge Base searcher\nTo get started, Please enter your search query:")
searchQuery = input()

# Searches for all the files in the directory
file_list = []
for root, directories, files in os.walk("."):
    for filename in files:
        file_list.append(os.path.join(root, filename))

matchingFiles = []

# Loops over all the files in the knowledge base and impliments a loading bar
for i in tqdm(file_list):

    # Gets the extension of the file so it knows which one to read
    extension = os.path.splitext(i)[1]
    if extension == ".pptx":
        test = readPPTX(i, searchQuery)
        if test != None:
            # Removes the root of the file path and adds the directory to the list to be printed later
            matchingFiles.append(test[6:])
    elif extension == ".pdf":
        test = readPDF(i, searchQuery)
        if test != None:
            matchingFiles.append(test[6:])
    elif extension == ".docx":
        test = readDOCX(i, searchQuery)
        if test != None:
            matchingFiles.append(test[6:])
    elif extension == ".xlsx":
        test = readXLSX(i, searchQuery)
        if test != None:
            matchingFiles.append(test[6:])

print("\nFiles containing the phrase/keyword: " + searchQuery + "\n")
for i in matchingFiles:
    print(i + "\n")