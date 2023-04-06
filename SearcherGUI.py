# Necessary imports for non-native file handling and directory traversal
import os
import docx
import openpyxl
import PyPDF2
import pptx
import tkinter as tk

# This function makes a tuple into all lowercase
def lowercaseTuple(tuple_val):
    if tuple_val is None:
        return None
    lowercase_vals = []
    for val in tuple_val:
        if val is not None:
            lowercase_vals.append(str(val).lower())
        else:
            lowercase_vals.append(None)
    return tuple(lowercase_vals)

# This function returns the path if it contains the desired keyword 
def readPDF(path, query):
    with open(path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
        if query.lower() in text.casefold() :
            return path

# This function returns the path if it contains the desired keyword 
def readDOCX(path, query):
    document = docx.Document(path)
    for paragraph in document.paragraphs:
        if query.lower() in paragraph.text.casefold():
            return path

# This function returns the path if it contains the desired keyword 
def readPPTX(path, query):
    prs = pptx.Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and query.lower() in shape.text.casefold():
                return path

# This function returns the path if it contains the desired keyword 
def readXLSX(path, query):
    workbook = openpyxl.load_workbook(path)
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            row = lowercaseTuple(row)
            if query.lower() in row:
                return path

# Main run
def run(event=0):
    searchQuery = query.get()
    results.delete(1.0, tk.END)

    # Searches for all the files in the directory
    file_list = []
    for root, directories, files in os.walk("."):
        for filename in files:
            file_list.append(os.path.join(root, filename))

    matchingFiles = []

    # Loops over all the files in the knowledge base and impliments a loading bar
    for i in file_list:

        # Gets the extension of the file so it knows which one to read
        extension = os.path.splitext(i)[1]

        if extension == ".pptx":
            test = readPPTX(i, searchQuery)
            if test != None:
                # Removes the root of the file path and adds the directory to the list to be printed later
                matchingFiles.append(test[6:])
        elif extension == ".pdf":
            test = readPDF(i, searchQuery)
            if test != None:
                matchingFiles.append(test[6:])
        elif extension == ".docx":
            test = readDOCX(i, searchQuery)
            if test != None:
                matchingFiles.append(test[6:])
        elif extension == ".xlsx":
            test = readXLSX(i, searchQuery)
            if test != None:
                matchingFiles.append(test[6:])

    resultLabel.config(text="Files containing the word: " + searchQuery)

    for i in matchingFiles:
        results.insert(1.0, i + "\n\n")

# Declaring window elements
window = tk.Tk()
window.title("Knowledge Base Searcher")
window.geometry("1000x600")

# Declaring elements for the window
queryLabel = tk.Label(window, text="Enter search query here: ")
query = tk.Entry(window)
searchButton = tk.Button(window, text="Search", command=run)
resultLabel = tk.Label(window, text="Files containing the word: ")
results = tk.Text(window, width=1000, height=400)


# Makes pressing the enter key run the same command as the search button
query.bind("<Return>", run)

# Packing all elements
queryLabel.pack(padx=5, pady=5)
query.pack(padx=5, pady=5)
searchButton.pack(padx=5, pady=5)
resultLabel.pack(padx=5, pady=5)
results.pack(padx=5, pady=5)

window.mainloop()